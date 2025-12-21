import sys
import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_ppt(pdf_path, output_ppt="output.pptx", dpi=1200):
    # --- Validate input file ---
    if not os.path.exists(pdf_path):
        print(f"❌ File not found: {pdf_path}")
        sys.exit(1)

    # --- Convert all PDF pages to images ---
    print(f"📄 Converting PDF pages to images from: {pdf_path}")
    pages = convert_from_path(pdf_path, dpi=dpi)
    total_pages = len(pages)
    print(f"✅ Found {total_pages} pages")

    # --- Prepare output directory ---
    temp_dir = "temp_pdf_images"
    os.makedirs(temp_dir, exist_ok=True)

    # --- Create PowerPoint presentation ---
    prs = Presentation()
    prs.slide_width = Inches(13.33)   # Default PPT width
    prs.slide_height = Inches(7.5)    # Default PPT height
    blank_slide_layout = prs.slide_layouts[6]

    # --- Convert each PDF page into a slide ---
    for i, page in enumerate(pages, start=1):
        img_path = os.path.join(temp_dir, f"page_{i}.png")
        page.save(img_path, "PNG")

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            img_path,
            left=Inches(0),
            top=Inches(0),
            width=Inches(13.33),
            height=Inches(7.5)
        )
        print(f"🖼️  Added page {i} as slide {i}")

    # --- Save presentation ---
    prs.save(output_ppt)
    print(f"\n🎉 Done! Saved PowerPoint file as: {output_ppt}")

    # --- (Optional) Cleanup temp images ---
    # import shutil
    # shutil.rmtree(temp_dir)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("❌ Usage: python pdf_to_ppt.py <pdf_path> [output_ppt]")
        sys.exit(1)

    pdf_path = sys.argv[1]
    output_ppt = sys.argv[2] if len(sys.argv) > 2 else pdf_path+".pptx"
    pdf_to_ppt(pdf_path, output_ppt)
